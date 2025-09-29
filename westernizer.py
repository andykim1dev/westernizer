from datetime import datetime
from PIL import Image
from openpyxl import load_workbook
from openpyxl.styles import Alignment, Border, PatternFill, Side
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt
import sys
import os
import os.path
import pandas as pd
import tkinter as tk
from tkinter import messagebox

def clean_marker(val):
    if pd.isna(val):
        return ""
    try:
        f = float(val)
        if f.is_integer():
            return str(int(f))
        return str(f)
    except:
        return str(val).strip()

os.environ["TK_SILENCE_DEPRECATION"] = "1"

# Declare directories, excel, and output ppt names
current_date_time = datetime.now().strftime('%y%m%d %H:%M')
current_date = datetime.now().strftime('%y.%m.%d')

if getattr(sys, 'frozen', False):
    base_dir = os.path.dirname(sys.executable)
else:
    base_dir = os.path.dirname(os.path.abspath(__file__))

wb_dir = os.path.join(base_dir, 'Insert WB')
marker_dir = os.path.join(base_dir, 'Insert MARKER')
ppt_dir = os.path.join(base_dir, 'Output PowerPoint')
excel_dir = os.path.join(base_dir, 'Output Excel')
excel_path = os.path.join(excel_dir, f'{current_date_time} WB Template.xlsx')
ppt_path = os.path.join(ppt_dir, f'{current_date_time} SUM.pptx')

# Extension for all images
extensions = ['.tif', '.jpg', '.png', '.jpeg']

# Slide dimensions
ppt = Presentation()
ppt.slide_width = Cm(33.87)
ppt.slide_height = Cm(19.05)
blank_slide_layout = ppt.slide_layouts[6]

# Create directories if they do not exist
for directory in [wb_dir, marker_dir, ppt_dir, excel_dir]:
    os.makedirs(directory, exist_ok=True)

# Prompt user to place Western Blot images in the specified directory
root = tk.Tk()
root.withdraw()
messagebox.showinfo("Insert WB Images", """Place Western Blot images in "Insert WB" and click OK when done.""")

# Check if there are any images in the directory
if not any(filename.lower().endswith(tuple(extensions)) for filename in os.listdir(wb_dir)):
    print(f"No images found in '{wb_dir}'. Please add images with extensions {extensions} and try again.")
    sys.exit()

# Make a list of all images in the directory that do not contain 'lad' in their name
image_files = [
    os.path.splitext(filename.lower())[0]
    for filename in os.listdir(wb_dir)
    if filename.lower().endswith(tuple(extensions))
    and "lad" not in os.path.splitext(filename)[0].lower()
]

# Create an information table
n = len(image_files)
data = {
    "Antibody Name": [name.upper() for name in image_files],
    "Volume (uL)": [""] * n,
    "Gel (%)": [""] * n,
    "Marker": [""] * n,
    "1' Antibody Catalog #": [""] * n,
    "1' Antibody Dilution (1 : x)": [""] * n,
    "": [""] * n,
    "Conditions": [""] * n
}
df = pd.DataFrame(data)
df.to_excel(excel_path, index=False)

# Reformat the excel file to widen the columns
wb = load_workbook(excel_path)
ws = wb.active
for col in ws.columns:
    col_letter = col[0].column_letter
    ws.column_dimensions[col_letter].width = 25

# Reformat the excel to highlight Columns A and D
highlight = PatternFill(start_color='FFFF00', end_color='FFFF00', fill_type='solid')
for cell in ws['A']:
    cell.fill = highlight
for cell in ws['D']:
    cell.fill = highlight
for cell in ws['H']:
    cell.fill = highlight

# Reformat the excel to add borders
thin_border = Border(left=Side(style='thin'),
                     right=Side(style='thin'), top=Side(style='thin'),
                     bottom=Side(style='thin'))
for row in ws.iter_rows():
    for cell in row:
        cell.border = thin_border

# Reformat the excel to center align text
center_alignment = Alignment(horizontal='center', vertical='center')
for row in ws.iter_rows():
    for cell in row:
        cell.alignment = center_alignment

# Save the reformatted Excel file
wb.save(excel_path)

# Open the Excel file for user input
os.system(f'open "{excel_path}"')

# Wait for user to fill out the Excel file
messagebox.showinfo("Fill Excel", f"Fill out 'WB Template.xlsx' and SAVE it. Click OK when done.")

# Reread the Excel file after user input
df = pd.read_excel(excel_path)

# Get 'Conditions' column, drop empty or NaN entries, convert to list
conditions = df["Conditions"].dropna()
if conditions.empty:
    print("No conditions found in the Excel file. Please fill out the 'Conditions' column and try again.")
    sys.exit()
conditions = conditions.astype(str).str.strip()
conditions = conditions[conditions != ""]

conditions = conditions.tolist()

# Use the updated DataFrame data (not the original lists)
antibody_names = df["Antibody Name"].dropna().tolist()
volumes = df["Volume (uL)"].tolist()
gel_percentages = df["Gel (%)"].tolist()
markers = df["Marker"].tolist()
catalog_numbers = df["1' Antibody Catalog #"].tolist()
dilutions = df["1' Antibody Dilution (1 : x)"].tolist()

for i in range(len(antibody_names)):
    print(antibody_names[i])

# Create a dictionary for antibody names and their corresponding marker
protein_dict = {
    str(antibody_names[i]).strip().lower(): clean_marker(markers[i])
    for i in range(len(antibody_names))
}

# Print the protein dictionary
for key, value in protein_dict.items():
    print(f'{key}: {value}')

# Get conditions from user input
conditions_count = len(conditions)

target_height = 19.05

# Parameters
label_top = Cm(15 * 0.7)
line_y = label_top - Cm(0.1 * 0.7)
line_width = Cm(1.1 * 0.7)
line_spacing = Cm(1.33 * 0.7)
label_width = Cm(1.5 * 0.7)
label_height = Cm(0.3 * 0.7)

slide_2 = None

# Center subtitle placeholder
slide_1 = ppt.slides.add_slide(blank_slide_layout)

# Create text box for title
left_1 = Cm(0)
top_1 = Cm(5.7)
width_1 = Cm(33.87)
height_1 = Cm(5)
textbox_1 = slide_1.shapes.add_textbox(left_1, top_1, width_1, height_1)
text_frame = textbox_1.text_frame
text_frame.text = "WB Summary: " + ", ".join([k.upper() for k in protein_dict.keys()])
text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Horizontal center
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
text_frame.paragraphs[0].font.size = Pt(32)
text_frame.paragraphs[0].font.bold = True
text_frame.word_wrap = True
text_frame.paragraphs[0].font.name = 'Arial'
text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

# Create text box for subtitle
left_2 = Cm(0)
top_2 = Cm(12)
width_2 = Cm(33.87)
height_2 = Cm(1)
textbox_2 = slide_1.shapes.add_textbox(left_2, top_2, width_2, height_2)
text_frame_2 = textbox_2.text_frame
text_frame_2.text = current_date
text_frame_2.paragraphs[0].alignment = PP_ALIGN.CENTER  # Horizontal center
text_frame_2.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
text_frame_2.paragraphs[0].font.size = Pt(16)
text_frame_2.paragraphs[0].font.bold = True
text_frame_2.word_wrap = True
text_frame_2.paragraphs[0].font.name = 'Arial'
text_frame_2.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

for filename in sorted(os.listdir(wb_dir)):
    if filename.lower().endswith(tuple(extensions)):
        image_path = os.path.join(wb_dir, filename)
        try:
            with Image.open(image_path) as img:
                width_px, height_px = img.size
        except Exception as e:
            print(f"Could not open image: {filename} — {e}")
            continue
        aspect_ratio = width_px / height_px

        # Calculate dimensions in cm
        height_cm = target_height
        width_cm = target_height * aspect_ratio

        # Convert to EMU
        width = Cm(width_cm)
        height = Cm(height_cm)

        # Add slide and image
        slide = ppt.slides.add_slide(blank_slide_layout)
        left = (ppt.slide_width - width) / 2
        top = (ppt.slide_height - height) / 2
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)

        # Add protein marker image if it exists
        marker_found = False
        marker_path = None
        filename_lower = filename.lower()
        for antibody_name in protein_dict.keys():
            if antibody_name in filename_lower:
                marker_name = protein_dict[antibody_name]
                if marker_name:  # Check if marker name is not empty
                    for ext in extensions:
                        candidate = os.path.join(marker_dir, marker_name + ext)
                        if os.path.exists(candidate):
                            marker_path = candidate
                            marker_found = True
                            break
                if marker_found:
                    break

        if not marker_found:
            print(f"No matching marker found for {filename}. Skipping marker image.")

        if marker_found and marker_path and os.path.exists(marker_path):
            try:
                with Image.open(marker_path) as marker_img:
                    marker_width_px, marker_height_px = marker_img.size
            except Exception as e:
                print(f"Could not open marker image: {marker_path} — {e}")
                continue
            marker_aspect_ratio = marker_width_px / marker_height_px

            # Calculate dimensions in cm
            marker_height_cm = 8.61
            marker_width_cm = marker_height_cm * marker_aspect_ratio

            # Convert to EMU
            marker_width = Cm(marker_width_cm)
            marker_height = Cm(marker_height_cm)

            # Add the marker image
            left = ppt.slide_width - marker_width
            top = ppt.slide_height - marker_height
            slide.shapes.add_picture(marker_path, left, top, width=marker_width, height=marker_height)
        else:
            print(f"Marker image not found for {filename}: {marker_path}")

        if slide_2 is None:
            slide_2 = slide

            # Add condition labels to the first slide only
            for i, condition in enumerate(conditions):
                x = Cm(16.935 - (conditions_count * line_spacing.cm / 2) + i * line_spacing.cm)

                # Connector line
                connector = slide_2.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    begin_x=x,
                    begin_y=line_y,
                    end_x=x + line_width,
                    end_y=line_y
                )
                line_format = connector.line
                line_format.width = Pt(2)
                line_format.fill.solid()
                line_format.fill.fore_color.rgb = RGBColor(0, 0, 0)  # black
                connector.shadow.inherit = False

                # Add textbox
                textbox = slide_2.shapes.add_textbox(x - (abs(line_width - label_width) / 2), label_top, label_width, label_height)
                text_frame = textbox.text_frame
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                if pd.notna(condition):
                    text_val = str(condition)
                    if text_val.endswith(".0"):
                        text_val = text_val[:-2]  # Remove trailing '.0'
                else:
                    text_val = ""
                run.text = text_val
                run.text = text_val
                run.font.size = Pt(6)
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # Add textbox on the left top corner
        left = Cm(0)
        top = Cm(0)
        width = Cm(24)
        height = Cm(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = os.path.splitext(filename)[0]
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # Add marker line
        textbox = slide.shapes.add_textbox(left, top + Cm(10), Cm(3), Cm(1))
        text_frame = textbox.text_frame
        text_frame.text = 'kDa'
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT  # Horizontal center
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            begin_x=left + Cm(3),
            begin_y=top + Cm(10.5),
            end_x=left + Cm(4),
            end_y=top + Cm(10.5)
        )
        line_format = connector.line
        line_format.width = Pt(2)
        line_format.fill.solid()
        line_format.fill.fore_color.rgb = RGBColor(0, 0, 0)
        connector.shadow.inherit = False

# Save presentation
ppt.save(ppt_path)
print(f"Saved: {ppt_path}")

messagebox.showinfo("Opening Completed PPT", """Presentation is completed and saved in "Output PPT".""")

os.system(f'open "{ppt_path}"')
